import os
from pptx import Presentation
from deep_translator import GoogleTranslator
import uuid
import json

UPLOAD_FOLDER = "uploads"
TRADUZIDOS_FOLDER = "traduzidos"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(TRADUZIDOS_FOLDER, exist_ok=True)

def traduzir_texto(texto, tradutor):
    if not texto or len(texto.strip()) == 0:
        return texto
    try:
        traduzido = tradutor.translate(texto)
        return traduzido if traduzido != texto else texto
    except Exception:
        return texto

def processar_traducao(file):
    """
    Processa a tradução do arquivo PowerPoint (.pptx)
    """
    # Salvar o arquivo enviado
    file_id = str(uuid.uuid4())
    caminho_arquivo = os.path.join(UPLOAD_FOLDER, f"{file_id}.pptx")
    file.save(caminho_arquivo)
    
    # Carregar apresentação
    prs = Presentation(caminho_arquivo)
    tradutor = GoogleTranslator(source="en", target="pt")
    log_traducoes = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texto_original = run.text
                    texto_traduzido = traduzir_texto(texto_original, tradutor)
                    if texto_original != texto_traduzido:
                        run.text = texto_traduzido
                        log_traducoes.append({
                            "texto_original": texto_original,
                            "texto_traduzido": texto_traduzido
                        })

    # Salvar apresentação traduzida
    caminho_saida = os.path.join(TRADUZIDOS_FOLDER, f"traduzido_{file_id}.pptx")
    prs.save(caminho_saida)

    # Salvar log
    caminho_log = os.path.join(TRADUZIDOS_FOLDER, f"log_{file_id}.json")
    with open(caminho_log, "w", encoding="utf-8") as log_file:
        json.dump(log_traducoes, log_file, ensure_ascii=False, indent=4)

    return caminho_saida, log_traducoes