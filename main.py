import os
import sys
import json
from docx import Document
import pdfplumber
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from deep_translator import GoogleTranslator
from glob import glob
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile 
import shutil
from app import create_app

FORMATOS_SUPORTADOS = [".docx", ".pdf", ".txt", ".pptx"]
IDIOMAS_SUPORTADOS = ["en", "es", "fr", "de", "it", "pt"]

def executar_modo_batch(config_path):
    try:
        with open(config_path, 'r', encoding='utf-8') as config_file:
            config_data = json.load(config_file)

        blocos = config_data.get("blocos", [])

        for bloco in blocos:
            origem = bloco.get("idioma_origem", "auto")
            destino = bloco.get("idioma_destino")
            diretoria = bloco.get("diretorio")

            if not destino or not diretoria:
                print(f"[ERRO] Bloco inválido: {bloco}")
                continue

            if destino not in IDIOMAS_SUPORTADOS:
                print(f"[AVISO] Idioma destino '{destino}' não é suportado.")
                continue

            if not os.path.isdir(diretoria):
                print(f"[ERRO] Diretório não encontrado: {diretoria}")
                continue

            print(f"\n🔄 Processando diretoria: {diretoria}")
            arquivos = glob(os.path.join(diretoria, "*"))
            arquivos_validos = [arq for arq in arquivos if os.path.splitext(arq)[1].lower() in FORMATOS_SUPORTADOS]

            for file_path in arquivos_validos:
                try:
                    print(f"📄 Traduzindo: {os.path.basename(file_path)} → {destino.upper()}")
                    paragraphs = carregar_documento(file_path)

                    translator = GoogleTranslator(source=origem, target=destino)
                    translated_paragraphs = []
                    for para in paragraphs:
                        if para.strip():
                            try:
                                translated_text = translator.translate(para.strip())
                                translated_paragraphs.append(translated_text)
                            except Exception as e:
                                print(f"[AVISO] Erro ao traduzir parágrafo: {e}")
                                translated_paragraphs.append(para)

                    base_filename = os.path.basename(file_path)
                    name, ext = os.path.splitext(base_filename)
                    novo_nome = f"{name}_{destino}{ext}"
                    output_path = os.path.join(diretoria, novo_nome)

                    salvar_documento_sem_dialog(file_path, translated_paragraphs, output_path, destino)  
                    print(f" Salvo: {output_path}")

                except Exception as e:
                    print(f"[ERRO] Falha ao processar {file_path}: {e}")

        print("\n Processamento concluído.")

    except Exception as e:
        print(f"[ERRO] Falha ao carregar arquivo de configuração: {e}")

def salvar_documento_sem_dialog(file_path, translated_paragraphs, output_path, destino):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".docx":

        doc = Document(file_path)
        index = 0

        
        for para in doc.paragraphs:
            if not para.text.strip():
                continue

            
            style = para.style
            alignment = para.paragraph_format.alignment
            
            
            runs_info = []
            for run in para.runs:
                if run.text.strip():
                    runs_info.append({
                        'text': run.text,
                        'bold': run.bold,
                        'italic': run.italic,
                        'underline': run.underline,
                        'font': run.font.name,
                        'size': run.font.size,
                        'color': run.font.color.rgb if hasattr(run.font.color, 'rgb') else None,
                        'has_image': bool('graphicData' in run._r.xml or 'picture' in run._r.xml)
                    })

           
            if runs_info and index < len(translated_paragraphs):
                
                xml_original = para._p.xml
                para.clear()
                
                
                for run_info in runs_info:
                    if run_info['has_image']:
                        
                        run_xml = run_info.get('xml', '')
                        if run_xml:
                            new_run = para.add_run()
                            new_run._r.append(run_xml)
                    else:
                       
                        new_run = para.add_run(translated_paragraphs[index])
                        new_run.bold = run_info['bold']
                        new_run.italic = run_info['italic']
                        new_run.underline = run_info['underline']
                        new_run.font.name = run_info['font']
                        if run_info['size']:
                            new_run.font.size = run_info['size']
                        if run_info['color']:
                            new_run.font.color.rgb = run_info['color']
                
                
                if any(not r['has_image'] for r in runs_info):
                    index += 1

                
                para.style = style
                para.paragraph_format.alignment = alignment

       
        doc.save(output_path)


    elif file_extension == ".pdf":
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        margin = 72
        y_position = height - margin

        for para in translated_paragraphs:
            words = para.split()
            current_line = []
            for word in words:
                current_line.append(word)
                if len(' '.join(current_line)) > 80:
                    c.drawString(margin, y_position, ' '.join(current_line[:-1]))
                    y_position -= 15
                    current_line = [word]
                    if y_position < margin:
                        c.showPage()
                        y_position = height - margin
            if current_line:
                c.drawString(margin, y_position, ' '.join(current_line))
                y_position -= 20
            if y_position < margin:
                c.showPage()
                y_position = height - margin

        c.save()

    elif file_extension == ".txt":
        with open(output_path, "w", encoding="utf-8") as file:
            file.writelines(translated_paragraphs)

    elif file_extension == ".pptx":
        traduzir_pptx_mantendo_formatacao(file_path, output_path, destino)

    else:
        raise ValueError("Formato de arquivo não suportado para salvar.")

def carregar_documento(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".docx":
        doc = Document(file_path)
        return [para.text for para in doc.paragraphs]

    elif file_extension == ".pdf":
        paragraphs = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    lines = text.split('\n')
                    paragraphs.extend([line for line in lines if line.strip()])
        return paragraphs

    elif file_extension == ".txt":
        with open(file_path, "r", encoding="utf-8") as file:
            return file.readlines()

    elif file_extension == ".pptx":
        return carregar_pptx(file_path)  

    else:
        raise ValueError("Formato de arquivo não suportado.")

def carregar_pptx(file_path):
    prs = Presentation(file_path)
    paragraphs = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragrafo in shape.text_frame.paragraphs:
                    for run in paragrafo.runs:
                        texto = run.text.strip()
                        if texto:
                            paragraphs.append(texto)
    return paragraphs

def salvar_imagem_temporaria(imagem_bytes, extension):
    """Salvar imagens temporária"""
    temp_dir = tempfile.mkdtemp()
    temp_path = os.path.join(temp_dir, f'temp_image.{extension}')
    
    with open(temp_path, 'wb') as f:
        f.write(imagem_bytes)
    
    return temp_path, temp_dir

def traduzir_pptx_mantendo_formatacao(caminho_entrada, caminho_saida, idioma_destino='pt'):
    prs = Presentation(caminho_entrada)
    tradutor = GoogleTranslator(source='auto', target=idioma_destino)
    temp_dirs = []  

    try:
        for slide in prs.slides:
            
            imagens_info = {}
            
            
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    imagens_info[shape.shape_id] = {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'zorder': shape.element.get_or_add_ln().attrib.get('z-order', 0),
                        'rotation': shape.rotation
                    }
                                       
                    if hasattr(shape, 'image'):
                        temp_path, temp_dir = salvar_imagem_temporaria(shape.image.blob, shape.image.ext)
                        imagens_info[shape.shape_id]['temp_path'] = temp_path
                        temp_dirs.append(temp_dir)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        alignment = paragraph.alignment
                        space_before = paragraph.space_before
                        space_after = paragraph.space_after
                        
                        for run in paragraph.runs:
                            font_name = run.font.name
                            font_size = run.font.size
                            bold = run.font.bold
                            italic = run.font.italic
                            underline = run.font.underline
                            
                            texto_original = run.text.strip()
                            if texto_original:
                                try:                                  
                                    run.text = tradutor.translate(texto_original)
                                    run.font.name = font_name
                                    run.font.size = font_size
                                    run.font.bold = bold
                                    run.font.italic = italic
                                    run.font.underline = underline
                                except Exception as e:
                                    print(f"Erro ao traduzir '{texto_original}': {e}")                                           
                        paragraph.alignment = alignment
                        paragraph.space_before = space_before
                        paragraph.space_after = space_after

 
            for shape_id, info in imagens_info.items():
                try:
                
                    picture = slide.shapes.add_picture(
                        info['temp_path'],
                        info['left'],
                        info['top'],
                        width=info['width'],
                        height=info['height']
                    )
                                   
                    picture.rotation = info['rotation']
                    if hasattr(picture.element, 'get_or_add_ln'):
                        picture.element.get_or_add_ln().attrib['z-order'] = str(info['zorder'])
                except Exception as e:
                    print(f"Erro ao restaurar imagem: {e}")

        
        prs.save(caminho_saida)

    finally:
        for temp_dir in temp_dirs:
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"Erro ao remover diretório temporário: {e}")


if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] == "--config":
        if len(sys.argv) < 3:
            print("Erro: Caminho para o arquivo de configuração não fornecido.")
            sys.exit(1)
        caminho_config = sys.argv[2]
        executar_modo_batch(caminho_config)
    else:
        print("Modo gráfico desativado. Use com: python main.py --config caminho_para_config.json")

# Criar a aplicação Flask
app = create_app()

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)