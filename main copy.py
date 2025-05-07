import os
import sys
import json
from docx import Document
import pdfplumber
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from deep_translator import GoogleTranslator
from glob import glob
from pptx import Presentation
from pptx.util import Inches


FORMATOS_SUPORTADOS = [".docx", ".pdf", ".txt", ".pptx"]
IDIOMAS_SUPORTADOS = ["en", "es", "fr", "de", "it", "pt"]

def executar_modo_batch(config_path):
    try:
        with open(config_path, 'r', encoding='utf-8') as config_file:
            config_data = json.load(config_file)

        blocos = config_data.get("blocos", [])

        for bloco in blocos:
            origem = bloco.get("idioma_origem", "auto")
            destino = bloco.get("idioma_destino")
            diretoria = bloco.get("diretorio")

            if not destino or not diretoria:
                print(f"[ERRO] Bloco inválido: {bloco}")
                continue

            if destino not in IDIOMAS_SUPORTADOS:
                print(f"[AVISO] Idioma destino '{destino}' não é suportado.")
                continue

            if not os.path.isdir(diretoria):
                print(f"[ERRO] Diretório não encontrado: {diretoria}")
                continue

            print(f"\n🔄 Processando diretoria: {diretoria}")
            arquivos = glob(os.path.join(diretoria, "*"))
            arquivos_validos = [arq for arq in arquivos if os.path.splitext(arq)[1].lower() in FORMATOS_SUPORTADOS]

            for file_path in arquivos_validos:
                try:
                    print(f"📄 Traduzindo: {os.path.basename(file_path)} → {destino.upper()}")
                    paragraphs = carregar_documento(file_path)

                    translator = GoogleTranslator(source=origem, target=destino)
                    translated_paragraphs = []
                    for para in paragraphs:
                        if para.strip():
                            try:
                                translated_text = translator.translate(para.strip())
                                translated_paragraphs.append(translated_text)
                            except Exception as e:
                                print(f"[AVISO] Erro ao traduzir parágrafo: {e}")
                                translated_paragraphs.append(para)

                    base_filename = os.path.basename(file_path)
                    name, ext = os.path.splitext(base_filename)
                    novo_nome = f"{name}_{destino}{ext}"
                    output_path = os.path.join(diretoria, novo_nome)

                    salvar_documento_sem_dialog(file_path, translated_paragraphs, output_path, destino)  # Agora passando o 'destino'
                    print(f"✅ Salvo: {output_path}")

                except Exception as e:
                    print(f"[ERRO] Falha ao processar {file_path}: {e}")

        print("\n🎉 Processamento concluído.")

    except Exception as e:
        print(f"[ERRO] Falha ao carregar arquivo de configuração: {e}")

def salvar_documento_sem_dialog(file_path, translated_paragraphs, output_path, destino):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".docx":
        doc = Document(file_path)
        index = 0
        for para in doc.paragraphs:
            if index < len(translated_paragraphs) and para.text.strip():
                para.text = translated_paragraphs[index]
                index += 1
        doc.save(output_path)

    elif file_extension == ".pdf":
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        margin = 72
        y_position = height - margin

        for para in translated_paragraphs:
            words = para.split()
            current_line = []
            for word in words:
                current_line.append(word)
                if len(' '.join(current_line)) > 80:
                    c.drawString(margin, y_position, ' '.join(current_line[:-1]))
                    y_position -= 15
                    current_line = [word]
                    if y_position < margin:
                        c.showPage()
                        y_position = height - margin
            if current_line:
                c.drawString(margin, y_position, ' '.join(current_line))
                y_position -= 20
            if y_position < margin:
                c.showPage()
                y_position = height - margin

        c.save()

    elif file_extension == ".txt":
        with open(output_path, "w", encoding="utf-8") as file:
            file.writelines(translated_paragraphs)

    elif file_extension == ".pptx":
        traduzir_pptx_mantendo_formatacao(file_path, output_path, destino)

    else:
        raise ValueError("Formato de arquivo não suportado para salvar.")



def carregar_documento(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".docx":
        doc = Document(file_path)
        return [para.text for para in doc.paragraphs]

    elif file_extension == ".pdf":
        paragraphs = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    lines = text.split('\n')
                    paragraphs.extend([line for line in lines if line.strip()])
        return paragraphs

    elif file_extension == ".txt":
        with open(file_path, "r", encoding="utf-8") as file:
            return file.readlines()

    elif file_extension == ".pptx":
        return carregar_pptx(file_path)  # Agora chama a função correta para carregar o PPTX

    else:
        raise ValueError("Formato de arquivo não suportado.")
  
def carregar_pptx(file_path):
    prs = Presentation(file_path)
    paragraphs = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragrafo in shape.text_frame.paragraphs:
                    for run in paragrafo.runs:
                        texto = run.text.strip()
                        if texto:
                            paragraphs.append(texto)
    return paragraphs
    
def traduzir_pptx_mantendo_formatacao(caminho_entrada, caminho_saida, idioma_destino='pt'):
    prs = Presentation(caminho_entrada)
    tradutor = GoogleTranslator(source='auto', target=idioma_destino)

    for slide in prs.slides:
        for shape in slide.shapes:
            # Se o shape tem texto
            if shape.has_text_frame:
                for paragrafo in shape.text_frame.paragraphs:
                    for run in paragrafo.runs:
                        texto_original = run.text.strip()
                        if texto_original:
                            try:
                                # Traduzir o texto
                                run.text = tradutor.translate(texto_original)
                            except Exception as e:
                                print(f"Erro ao traduzir '{texto_original}': {e}")
            
            # Se o shape tem imagem, preservar a imagem
            if shape.shape_type == 13:  # Tipo 13 é imagem
                imagem = shape.image
                imagem_bytes = imagem.blob
                imagem_extension = imagem.ext

                # Inserir imagem novamente
                imagem_path = os.path.join('temp_image.' + imagem_extension)  # Salvar como arquivo temporário
                with open(imagem_path, 'wb') as img_file:
                    img_file.write(imagem_bytes)

                # Adicionar imagem preservada ao slide
                left = Inches(1)
                top = Inches(1)
                slide.shapes.add_picture(imagem_path, left, top)

    # Salvar o novo PowerPoint com tradução
    prs.save(caminho_saida)


if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] == "--config":
        if len(sys.argv) < 3:
            print("Erro: Caminho para o arquivo de configuração não fornecido.")
            sys.exit(1)
        caminho_config = sys.argv[2]
        executar_modo_batch(caminho_config)
    else:
        print("Modo gráfico desativado. Use com: python main.py --config caminho_para_config.json")
